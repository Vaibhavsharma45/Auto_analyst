"""
backend/analysis/ppt_generator.py
Auto PowerPoint generation from analysis data
"""
import pandas as pd
import numpy as np
import io
import os
import base64
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Dark theme colors
DARK_BG   = RGBColor(0x0D, 0x11, 0x17) if PPTX_AVAILABLE else None
ACCENT    = RGBColor(0x58, 0xA6, 0xFF) if PPTX_AVAILABLE else None
GREEN     = RGBColor(0x3F, 0xB9, 0x50) if PPTX_AVAILABLE else None
ORANGE    = RGBColor(0xD2, 0x99, 0x22) if PPTX_AVAILABLE else None
RED       = RGBColor(0xF8, 0x51, 0x49) if PPTX_AVAILABLE else None
WHITE     = RGBColor(0xE6, 0xED, 0xF3) if PPTX_AVAILABLE else None
MUTED     = RGBColor(0x8B, 0x94, 0x9E) if PPTX_AVAILABLE else None
PURPLE    = RGBColor(0xBC, 0x8C, 0xFF) if PPTX_AVAILABLE else None
CARD_BG   = RGBColor(0x1C, 0x23, 0x33) if PPTX_AVAILABLE else None


def _set_slide_bg(slide, color):
    from pptx.util import Pt
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or WHITE
    return txBox


def _add_rect(slide, x, y, w, h, color, radius=False):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


class PPTGenerator:
    def __init__(self, df: pd.DataFrame, analysis: dict, charts: dict, goal: dict = None, filename: str = "dataset"):
        self.df = df
        self.analysis = analysis
        self.charts = charts
        self.goal = goal or {}
        self.filename = filename
        self.date = datetime.now().strftime("%B %d, %Y")

    def generate(self, output_path: str) -> str:
        if not PPTX_AVAILABLE:
            return self._fallback(output_path)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        self._slide_cover(prs, blank_layout)
        self._slide_overview(prs, blank_layout)
        self._slide_data_quality(prs, blank_layout)
        self._slide_numeric_stats(prs, blank_layout)
        self._slide_charts(prs, blank_layout)
        self._slide_correlations(prs, blank_layout)
        self._slide_insights(prs, blank_layout)
        self._slide_recommendations(prs, blank_layout)
        self._slide_thank_you(prs, blank_layout)

        prs.save(output_path)
        return output_path

    def _slide_cover(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        # Accent bar top
        _add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
        # Title
        _add_textbox(slide, "⚡ DataMind Pro", 1, 1.2, 11, 1, size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _add_textbox(slide, "Automated Data Analysis Report", 1, 2.3, 11, 0.6, size=22, color=WHITE, align=PP_ALIGN.CENTER)
        # Divider
        _add_rect(slide, 3, 3.1, 7.33, 0.03, ACCENT)
        # Meta info
        meta = f"Dataset: {self.filename}   |   Generated: {self.date}"
        _add_textbox(slide, meta, 1, 3.3, 11, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
        shape_info = self.analysis.get("overview", {}).get("shape", {})
        info = f"{shape_info.get('rows',0):,} rows  ×  {shape_info.get('columns',0)} columns"
        _add_textbox(slide, info, 1, 3.9, 11, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        qr = self.analysis.get("quality_report", {})
        score = qr.get("quality_score", 0)
        score_color = GREEN if score >= 80 else ORANGE if score >= 60 else RED
        _add_textbox(slide, f"Quality Score: {score}/100", 1, 4.6, 11, 0.5, size=16, color=score_color, align=PP_ALIGN.CENTER)
        if self.goal.get("question"):
            _add_textbox(slide, f'🎯 "{self.goal["question"]}"', 1, 5.3, 11, 0.6, size=14, color=PURPLE, align=PP_ALIGN.CENTER)
        _add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)

    def _slide_overview(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, ACCENT)
        _add_textbox(slide, "Dataset Overview", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=ACCENT)

        ov = self.analysis.get("overview", {})
        sh = ov.get("shape", {})
        mis = ov.get("missing", {})
        dup = ov.get("duplicates", {})
        ct = ov.get("column_types", {})
        qr = self.analysis.get("quality_report", {})

        stats = [
            ("Rows", f"{sh.get('rows',0):,}", ACCENT),
            ("Columns", str(sh.get("columns",0)), PURPLE),
            ("Missing", f"{mis.get('total_missing_cells',0)}", ORANGE if mis.get('total_missing_cells',0) > 0 else GREEN),
            ("Duplicates", str(dup.get("duplicate_rows",0)), ORANGE if dup.get("duplicate_rows",0) > 0 else GREEN),
            ("Quality", f"{qr.get('quality_score',0)}/100", GREEN if qr.get("quality_score",0)>=80 else ORANGE),
            ("Numeric Cols", str(len(ct.get("numeric",[]))), ACCENT),
        ]

        for i, (label, val, color) in enumerate(stats):
            x = 0.5 + (i % 3) * 4.1
            y = 1.1 + (i // 3) * 1.6
            _add_rect(slide, x, y, 3.6, 1.3, CARD_BG)
            _add_textbox(slide, label, x+0.15, y+0.1, 3.3, 0.4, size=11, color=MUTED)
            _add_textbox(slide, val, x+0.15, y+0.45, 3.3, 0.7, size=30, bold=True, color=color)

        # Column names
        _add_textbox(slide, "Columns: " + ", ".join(self.df.columns.tolist()), 0.5, 4.7, 12.3, 0.5, size=11, color=MUTED)

    def _slide_data_quality(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, ORANGE)
        _add_textbox(slide, "Data Quality Report", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=ORANGE)

        qr = self.analysis.get("quality_report", {})
        issues = qr.get("issues", [])[:6]
        recs = qr.get("recommendations", [])[:4]

        _add_textbox(slide, f"Quality Score: {qr.get('quality_score',0)}/100", 0.5, 1.0, 5, 0.5, size=20, bold=True,
                     color=GREEN if qr.get("quality_score",0)>=80 else ORANGE)

        y = 1.7
        sev_colors = {"critical": RED, "warning": ORANGE, "info": ACCENT}
        for issue in issues:
            _add_textbox(slide, f"• {issue['message'][:80]}", 0.5, y, 7.5, 0.4, size=11,
                         color=sev_colors.get(issue["severity"], WHITE))
            y += 0.4

        if recs:
            _add_textbox(slide, "Recommendations:", 8.5, 1.0, 4.5, 0.4, size=14, bold=True, color=GREEN)
            ry = 1.6
            for rec in recs:
                _add_textbox(slide, f"→ {rec[:50]}", 8.5, ry, 4.5, 0.5, size=11, color=WHITE)
                ry += 0.5

    def _slide_numeric_stats(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, ACCENT)
        _add_textbox(slide, "Statistical Analysis", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=ACCENT)

        ns = self.analysis.get("numeric_stats", {})
        headers = ["Column", "Mean", "Median", "Std", "Min", "Max", "Outliers"]
        col_w = [2.5, 1.4, 1.4, 1.4, 1.4, 1.4, 1.3]
        x_starts = [0.3]
        for w in col_w[:-1]:
            x_starts.append(x_starts[-1] + w)

        y = 1.0
        # Header row
        _add_rect(slide, 0.3, y, 12.7, 0.4, CARD_BG)
        for i, (h, x) in enumerate(zip(headers, x_starts)):
            _add_textbox(slide, h, x+0.05, y+0.05, col_w[i]-0.1, 0.3, size=10, bold=True, color=ACCENT)
        y += 0.45

        for col, s in list(ns.items())[:7]:
            _add_rect(slide, 0.3, y, 12.7, 0.38, RGBColor(0x16,0x1B,0x22))
            vals = [col[:18], str(s.get("mean","")), str(s.get("median","")), str(s.get("std","")),
                    str(s.get("min","")), str(s.get("max","")), str(s.get("outliers",{}).get("count",0))]
            for i, (v, x) in enumerate(zip(vals, x_starts)):
                _add_textbox(slide, str(v)[:16], x+0.05, y+0.04, col_w[i]-0.1, 0.3, size=9,
                             color=ACCENT if i==0 else WHITE)
            y += 0.4

    def _slide_charts(self, prs, layout):
        """Add chart images to slides"""
        chart_keys = ["distribution_dashboard", "correlation_heatmap", "boxplot_dashboard", "scatter_matrix"]
        chart_titles = ["Distributions", "Correlation Heatmap", "Box Plots", "Scatter Matrix"]

        for key, title in zip(chart_keys, chart_titles):
            if key not in self.charts:
                continue
            chart_val = self.charts[key]
            if not isinstance(chart_val, str) or not chart_val.startswith("data:image/png;base64,"):
                continue

            slide = prs.slides.add_slide(layout)
            _set_slide_bg(slide, DARK_BG)
            _add_rect(slide, 0, 0, 13.33, 0.06, PURPLE)
            _add_textbox(slide, f"📊 {title}", 0.5, 0.15, 12, 0.55, size=24, bold=True, color=PURPLE)

            try:
                img_data = base64.b64decode(chart_val.split(",")[1])
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.3))
            except Exception:
                _add_textbox(slide, "Chart not available", 0.5, 3, 12, 1, size=16, color=MUTED, align=PP_ALIGN.CENTER)

    def _slide_correlations(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, PURPLE)
        _add_textbox(slide, "Correlation Analysis", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=PURPLE)

        corr = self.analysis.get("correlations", {})
        strong = corr.get("strong_pairs", [])

        if not strong:
            _add_textbox(slide, "No strong correlations found (|r| < 0.5)", 0.5, 2, 12, 1, size=18, color=MUTED, align=PP_ALIGN.CENTER)
            return

        _add_textbox(slide, f"Found {len(strong)} strong correlations (|r| ≥ 0.5):", 0.5, 1.0, 12, 0.4, size=14, color=MUTED)
        y = 1.6
        for pair in strong[:8]:
            color = GREEN if pair["pearson_r"] > 0 else RED
            bar_w = abs(pair["pearson_r"]) * 5
            _add_rect(slide, 0.5, y, bar_w, 0.35, RGBColor(0x1C,0x23,0x33))
            _add_rect(slide, 0.5, y, bar_w, 0.35, color)
            text = f"{pair['col1']} ↔ {pair['col2']}   r = {pair['pearson_r']}   ({pair['strength']})"
            _add_textbox(slide, text, 6.0, y, 7, 0.35, size=11, color=WHITE)
            y += 0.5

    def _slide_insights(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, GREEN)
        _add_textbox(slide, "Key Insights", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=GREEN)

        ns = self.analysis.get("numeric_stats", {})
        qr = self.analysis.get("quality_report", {})
        y = 1.1

        insights = []
        for col, s in list(ns.items())[:4]:
            insights.append(f"📊 {col}: Mean={s.get('mean')}, Range=[{s.get('min')}, {s.get('max')}], Outliers={s.get('outliers',{}).get('count',0)}")
        for issue in qr.get("issues", [])[:3]:
            icon = "🔴" if issue["severity"]=="critical" else "⚠️" if issue["severity"]=="warning" else "ℹ️"
            insights.append(f"{icon} {issue['message'][:80]}")

        for ins in insights[:7]:
            _add_rect(slide, 0.5, y, 12.3, 0.5, CARD_BG)
            _add_textbox(slide, ins, 0.7, y+0.05, 11.8, 0.4, size=12, color=WHITE)
            y += 0.6

    def _slide_recommendations(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, ORANGE)
        _add_textbox(slide, "Recommendations & Next Steps", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=ORANGE)
        _add_textbox(slide, '"Insight without action = useless."', 0.5, 0.95, 12, 0.4, size=13, color=MUTED)

        qr = self.analysis.get("quality_report", {})
        recs = qr.get("recommendations", [])[:6]

        y = 1.6
        for i, rec in enumerate(recs):
            _add_rect(slide, 0.5, y, 0.4, 0.4, ORANGE)
            _add_textbox(slide, str(i+1), 0.55, y+0.02, 0.3, 0.35, size=14, bold=True, color=DARK_BG)
            _add_rect(slide, 1.1, y, 11.7, 0.4, CARD_BG)
            _add_textbox(slide, rec[:90], 1.25, y+0.04, 11.4, 0.32, size=11, color=WHITE)
            y += 0.55

    def _slide_thank_you(self, prs, layout):
        slide = prs.slides.add_slide(layout)
        _set_slide_bg(slide, DARK_BG)
        _add_rect(slide, 0, 0, 13.33, 0.06, ACCENT)
        _add_textbox(slide, "Thank You", 1, 2.0, 11, 1.2, size=52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _add_textbox(slide, "Generated by DataMind Pro", 1, 3.5, 11, 0.6, size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _add_textbox(slide, self.date, 1, 4.2, 11, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)
        _add_rect(slide, 0, 7.44, 13.33, 0.06, ACCENT)

    def _fallback(self, path: str) -> str:
        fallback_path = path.replace(".pptx", ".txt")
        with open(fallback_path, "w") as f:
            f.write("DataMind Pro — PPT\npip install python-pptx to enable PowerPoint generation\n")
        return fallback_path
